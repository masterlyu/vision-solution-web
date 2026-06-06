# -*- coding: utf-8 -*-
"""슬라이드 PNG들 → PPTX (각 슬라이드 = 풀블리드 이미지). 코스1과 동일 구조."""
import os, glob
from pptx import Presentation
from pptx.util import Inches

SHOTS = "/home/ubuntu/company/website/scripts/academy/_shots"
OUT = "/home/ubuntu/company/website/src/storage/academy/build-ai-part1-slides.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

pngs = sorted(glob.glob(os.path.join(SHOTS, "slide-*.png")))
assert pngs, "스크린샷 없음"
for p in pngs:
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(p, 0, 0, width=prs.slide_width, height=prs.slide_height)

prs.save(OUT)
print("PPTX 저장:", OUT, "/", len(pngs), "슬라이드 /", os.path.getsize(OUT)//1024, "KB")
