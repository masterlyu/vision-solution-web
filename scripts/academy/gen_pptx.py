# -*- coding: utf-8 -*-
"""buildai_part1_data.SLIDES → PPTX (편집 가능한 브랜드 덱)."""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

sys.path.insert(0, "/home/ubuntu/company/website/scripts/academy")
from buildai_part1_data import SLIDES, TOTAL, META

# ---- 브랜드 색 ----
INK      = RGBColor(0x0F, 0x17, 0x2A)   # 본문 잉크
MUTED    = RGBColor(0x64, 0x74, 0x8B)   # 보조 텍스트
PRIMARY  = RGBColor(0x0E, 0x9F, 0x9A)   # 비젼 teal
BLUE     = RGBColor(0x25, 0x63, 0xEB)   # 🟦 불변층
AMBER    = RGBColor(0xD9, 0x77, 0x06)   # 🟨 가변층
PAPER    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x0B, 0x12, 0x20)   # 타이틀 슬라이드 배경
LINE     = RGBColor(0xE2, 0xE8, 0xF0)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def _set(run, size, color, bold=False, font=None):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    if font:
        run.font.name = font


def _fill(slide, l, t, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def sub_color(sub):
    if "가변" in sub:
        return AMBER
    if "불변" in sub:
        return BLUE
    return PRIMARY


# ===================== 타이틀 슬라이드 =====================
s = prs.slides.add_slide(BLANK)
_fill(s, 0, 0, EMU_W, EMU_H, DARK)
_fill(s, 0, Inches(4.05), EMU_W, Inches(0.06), PRIMARY)

_, tf = _box(s, Inches(0.9), Inches(0.8), Inches(11.5), Inches(0.6))
r = tf.paragraphs[0].add_run(); r.text = "VISIONC · ENTERPRISE ACADEMY"
_set(r, 16, PRIMARY, bold=True)

_, tf = _box(s, Inches(0.9), Inches(1.45), Inches(11.5), Inches(0.6))
r = tf.paragraphs[0].add_run()
r.text = f"{META['course']}".upper() if False else META["course"]
_set(r, 18, RGBColor(0xCB, 0xD5, 0xE1), bold=False)

_, tf = _box(s, Inches(0.85), Inches(2.5), Inches(11.6), Inches(1.7))
p = tf.paragraphs[0]
r = p.add_run(); r.text = f"{META['part']} · {META['part_title']}"
_set(r, 30, PRIMARY, bold=True)
p2 = tf.add_paragraph()
r = p2.add_run(); r.text = META["subtitle"]
_set(r, 40, PAPER, bold=True)

# 메타 박스
meta_lines = [
    ("코스", META["course"]),
    ("강의", META["lessons"] + "  ·  " + META["duration"]),
    ("대상", META["audience"]),
    ("형식", META["format"]),
]
_, tf = _box(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(2.4))
for i, (k, v) in enumerate(meta_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6)
    rk = p.add_run(); rk.text = f"{k}   "
    _set(rk, 14, PRIMARY, bold=True)
    rv = p.add_run(); rv.text = v
    _set(rv, 14, RGBColor(0xE2, 0xE8, 0xF0))


# ===================== 콘텐츠 슬라이드 =====================
for sl in SLIDES:
    s = prs.slides.add_slide(BLANK)
    _fill(s, 0, 0, EMU_W, EMU_H, PAPER)
    # 상단 라인
    _fill(s, Inches(0.9), Inches(1.25), Inches(11.53), Pt(2.2), PRIMARY)

    # 카테고리 (좌상)
    _, tf = _box(s, Inches(0.85), Inches(0.55), Inches(8), Inches(0.5))
    r = tf.paragraphs[0].add_run(); r.text = sl["cat"]
    _set(r, 13, PRIMARY, bold=True)

    # 슬라이드 번호 (우상)
    _, tf = _box(s, Inches(9.4), Inches(0.55), Inches(3.0), Inches(0.5))
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r = tf.paragraphs[0].add_run()
    r.text = f"{META['part']} · SLIDE {sl['no']:02d} / {TOTAL}"
    _set(r, 11, MUTED, bold=True)

    # 헤드라인
    _, tf = _box(s, Inches(0.85), Inches(1.55), Inches(11.6), Inches(1.3))
    r = tf.paragraphs[0].add_run(); r.text = sl["headline"]
    _set(r, 33, INK, bold=True)

    # 부제 (불변/가변 색)
    top_bul = Inches(3.05)
    if sl.get("sub"):
        _, tf = _box(s, Inches(0.87), Inches(2.75), Inches(11.6), Inches(0.5))
        r = tf.paragraphs[0].add_run(); r.text = sl["sub"]
        _set(r, 15, sub_color(sl["sub"]), bold=True)
        top_bul = Inches(3.45)

    # 불릿
    _, tf = _box(s, Inches(0.9), top_bul, Inches(11.5), Inches(3.4))
    for i, b in enumerate(sl["bullets"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        rb = p.add_run(); rb.text = "▸  "
        _set(rb, 18, PRIMARY, bold=True)
        rt = p.add_run(); rt.text = b
        _set(rt, 18, INK, bold=False)

    # 풋터
    _fill(s, Inches(0.9), Inches(6.95), Inches(11.53), Pt(1), LINE)
    _, tf = _box(s, Inches(0.85), Inches(7.0), Inches(11.6), Inches(0.4))
    r = tf.paragraphs[0].add_run()
    r.text = f"VISIONC · ENTERPRISE ACADEMY    ·    {META['course']}    ·    {META['part']} {META['part_title']}"
    _set(r, 9, MUTED, bold=False)

out = "/home/ubuntu/company/website/src/storage/academy/build-ai-part1-slides.pptx"
prs.save(out)
print("PPTX 저장:", out, "/ 슬라이드", len(prs.slides._sldIdLst), "장")
